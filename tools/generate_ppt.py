import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define helper functions to modify existing slides
def replace_text_in_slide(slide, target, replacement):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if target in shape.text:
                # Replace in all paragraphs
                for paragraph in shape.text_frame.paragraphs:
                    if target in paragraph.text:
                        paragraph.text = paragraph.text.replace(target, replacement)

def delete_slides_from(prs, start_idx):
    id_list = prs.slides._sldIdLst
    while len(prs.slides) > start_idx:
        del id_list[start_idx]

def set_shape_font(shape, font_name='Calibri', font_size_pt=18, bold=False, color_rgb=None):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.size = Pt(font_size_pt)
            paragraph.font.bold = bold
            if color_rgb:
                paragraph.font.color.rgb = color_rgb

def main():
    prs_path = 'Phase -II PPT.pptx'
    if not os.path.exists(prs_path):
        print(f"Error: {prs_path} not found.")
        return
        
    prs = Presentation(prs_path)
    print(f"Loaded {prs_path}. Slide count: {len(prs.slides)}")
    
    # ----------------------------------------------------
    # Modify Slide 1: Cover slide
    # ----------------------------------------------------
    slide_1 = prs.slides[0]
    # Correct spelling mistake "Societey's" to "Society's"
    replace_text_in_slide(slide_1, "Societey’s", "Society's")
    
    # ----------------------------------------------------
    # Modify Slide 2: Project Presentation Title
    # ----------------------------------------------------
    slide_2 = prs.slides[1]
    # Replace Sign Language text with PhantomShield details
    for shape in slide_2.shapes:
        if shape.has_text_frame:
            text = shape.text
            if "SIGN LANGUAGE" in text:
                shape.text_frame.text = ""
                p = shape.text_frame.paragraphs[0]
                p.text = "Project Presentation on\n\"PHANTOMSHIELD: AI-ENABLED PHISHING DETECTION AND EXPLANATION SYSTEM\""
                p.font.name = 'Calibri'
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x0F, 0x2C, 0x59)
            elif "Coordinator" in text:
                shape.text_frame.text = ""
                p = shape.text_frame.paragraphs[0]
                p.text = "Project Guide: Dr. Rahul M. Dhokane"
                p.font.name = 'Calibri'
                p.font.size = Pt(18)
                p.font.bold = True
            elif "Guide" in text:
                shape.text_frame.text = ""
                p = shape.text_frame.paragraphs[0]
                p.text = "Department of Information Technology"
                p.font.name = 'Calibri'
                p.font.size = Pt(18)
                p.font.bold = True

    # ----------------------------------------------------
    # Modify Slide 3: Student Details
    # ----------------------------------------------------
    slide_3 = prs.slides[2]
    for shape in slide_3.shapes:
        if shape.has_text_frame:
            text = shape.text
            if "Piyush" in text or "Chavan" in text or "Rutik" in text:
                shape.text_frame.text = ""
                students = [
                    "Mr. Wakchaure Sanchit Sanjay (Seat no- B400270167)",
                    "Miss. Kale Jayshree Sandip (Seat no- B400270126)",
                    "Miss. Dange Shreya Rajesh (Seat no- B400270113)",
                    "Mr. Wakchaure Ganesh Shivaji (Seat no- B400270166)"
                ]
                for idx, student in enumerate(students):
                    if idx == 0:
                        p = shape.text_frame.paragraphs[0]
                    else:
                        p = shape.text_frame.add_paragraph()
                    p.text = student
                    p.font.name = 'Calibri'
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ----------------------------------------------------
    # Modify Slide 4: Table of Contents
    # ----------------------------------------------------
    slide_4 = prs.slides[3]
    for shape in slide_4.shapes:
        if shape.has_text_frame:
            text = shape.text
            if "Abstract" in text and "Objective" in text:
                shape.text_frame.text = ""
                contents = [
                    "Abstract & Introduction",
                    "Literature Survey & Research Gaps",
                    "Problem Definition & Objectives",
                    "Proposed System Architecture",
                    "Feature Extraction & Optimization (RSTHFS)",
                    "Ensemble Classifier (RF + XGBoost) & Heuristics",
                    "Explainable AI (XAI) using SHAP & Gemini AI",
                    "UML & Data Flow Diagrams (DFD)",
                    "System Implementation & Screenshots",
                    "Experimental Results & Future Work"
                ]
                for idx, item in enumerate(contents):
                    if idx == 0:
                        p = shape.text_frame.paragraphs[0]
                    else:
                        p = shape.text_frame.add_paragraph()
                    p.text = f"{idx+1}. {item}"
                    p.font.name = 'Calibri'
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ----------------------------------------------------
    # Delete the remaining slides (indices 4 onwards)
    # ----------------------------------------------------
    delete_slides_from(prs, 4)
    print(f"Deleted old slides. Current slide count: {len(prs.slides)}")
    
    # ----------------------------------------------------
    # Slide Layouts mapping
    # Layout 0: Title Slide
    # Layout 1: Title and Content
    # Layout 2: Two Content
    # Layout 3: Title Only
    # Layout 4: Blank
    # ----------------------------------------------------
    
    # Helper to add standard bullet slide
    def add_bullet_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        set_shape_font(slide.shapes.title, font_size_pt=32, bold=True, color_rgb=RGBColor(0x0F, 0x2C, 0x59))
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                
            p.text = bullet
            p.font.name = 'Calibri'
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            # Indent handles level indentation
            if bullet.strip().startswith("- "):
                p.text = bullet.replace("- ", "", 1)
                p.level = 1
            else:
                p.level = 0
        return slide

    # Helper to add centered image slide
    def add_image_slide(title, image_path, width_inches=8):
        slide = prs.slides.add_slide(prs.slide_layouts[3]) # Title Only
        slide.shapes.title.text = title
        set_shape_font(slide.shapes.title, font_size_pt=32, bold=True, color_rgb=RGBColor(0x0F, 0x2C, 0x59))
        
        slide_width_in = prs.slide_width / 914400.0
        width = Inches(width_inches)
        left = (prs.slide_width - width) / 2
        top = Inches(2.1)
        
        slide.shapes.add_picture(image_path, left, top, width=width)
        return slide

    # Helper to add side-by-side images slide
    def add_two_images_slide(title, img_1, img_2, label_1="", label_2=""):
        slide = prs.slides.add_slide(prs.slide_layouts[3]) # Title Only
        slide.shapes.title.text = title
        set_shape_font(slide.shapes.title, font_size_pt=32, bold=True, color_rgb=RGBColor(0x0F, 0x2C, 0x59))
        
        slide_width_in = prs.slide_width / 914400.0
        if slide_width_in > 11: # Widescreen
            width = Inches(5.2)
            left_1 = Inches(1.0)
            left_2 = Inches(7.0)
        else: # Standard
            width = Inches(4.3)
            left_1 = Inches(0.5)
            left_2 = Inches(5.2)
            
        top = Inches(2.2)
        slide.shapes.add_picture(img_1, left_1, top, width=width)
        slide.shapes.add_picture(img_2, left_2, top, width=width)
        
        # Add labels if provided
        if label_1:
            txBox = slide.shapes.add_textbox(left_1, top + Inches(4.2), width, Inches(0.5))
            txBox.text_frame.text = label_1
            set_shape_font(txBox, font_size_pt=14, bold=True, color_rgb=RGBColor(0x33, 0x33, 0x33))
        if label_2:
            txBox = slide.shapes.add_textbox(left_2, top + Inches(4.2), width, Inches(0.5))
            txBox.text_frame.text = label_2
            set_shape_font(txBox, font_size_pt=14, bold=True, color_rgb=RGBColor(0x33, 0x33, 0x33))
            
        return slide

    # ----------------------------------------------------
    # Slide 5: Abstract
    # ----------------------------------------------------
    bullets_abstract = [
        "Problem Scope:",
        "- Phishing is a major cybersecurity threat across web, email, and mobile channels.",
        "- Users lack trust in warnings because traditional systems are opaque \"black-boxes\".",
        "Proposed System:",
        "- PhantomShield provides real-time phishing detection for URLs, emails, and SMS messages.",
        "- Powered by a Python FastAPI backend, a React frontend, and a Chrome extension plugin.",
        "Key Innovations:",
        "- Optimized Feature Space using Rough Set Theory (RSTHFS) for lightweight execution.",
        "- Stacking Ensemble Model combining Random Forest and XGBoost with a heuristics override layer.",
        "- Explainable AI (XAI) using SHAP values and Gemini API to generate user-friendly risk narratives."
    ]
    add_bullet_slide("Abstract", bullets_abstract)

    # ----------------------------------------------------
    # Slide 6: Introduction
    # ----------------------------------------------------
    bullets_intro = [
        "Phishing Attack Vector:",
        "- Exploits human trust and cognitive biases rather than software vulnerabilities.",
        "- Attackers mimic trusted brands, hide behind homograph domains, and use urgency language.",
        "Traditional Protection Methods:",
        "- Manual blacklists: Suffer from update delays, vulnerable to zero-hour attacks.",
        "- Signature & heuristic filters: Struggle to generalize, generate high false positives.",
        "Need for Explainable Security:",
        "- Binary classification is not enough; users need descriptive transparency.",
        "- Visual threat scores, attack-type labeling, and educational tips empower users to make safer choices."
    ]
    add_bullet_slide("Introduction", bullets_intro)

    # ----------------------------------------------------
    # Slide 7: Literature Survey
    # ----------------------------------------------------
    bullets_lit = [
        "1. Hybrid Machine Learning for URL-Based Detection (Karim et al.)",
        "- Proved the effectiveness of combining structural and lexical features in ML models.",
        "- Left a gap in explaining classification results to end-users.",
        "2. RSTHFS: Rough Set Theory Feature Selection (Setu et al.)",
        "- Introduced CDF-g ranking and rough set aggregation to optimize features.",
        "- Reduced feature redundancy by 69% and runtime by 61% without accuracy loss.",
        "3. Deep Learning via ResMLP (Remya et al.)",
        "- Proposed residual skip connections to solve vanishing gradients in deep networks.",
        "4. Explainable AI via SHAP (Banerjee & Mehta)",
        "- Validated SHAP values to explain individual classification decisions in cyber security."
    ]
    add_bullet_slide("Literature Survey", bullets_lit)

    # ----------------------------------------------------
    # Slide 8: Research Gaps
    # ----------------------------------------------------
    bullets_gaps = [
        "High Feature Redundancy:",
        "- Existing models extract up to 100 features, leading to high computational overhead.",
        "- This makes real-time deployment (like browser plugins) slow and impractical.",
        "Opaque ML Decisions (Black-Box Models):",
        "- Deep learning or ensemble models provide high accuracy but no reasoning.",
        "- Users ignore warnings they do not understand.",
        "Poor Zero-Day Detection:",
        "- Heavy reliance on blacklists leaves a 30% exposure gap for new domains.",
        "PhantomShield Resolution:",
        "- Applies RSTHFS to reduce latency to <50ms, enabling extension integration.",
        "- Fuses SHAP and Gemini API to generate plain-text explanations of threats."
    ]
    add_bullet_slide("Research Gap Analysis", bullets_gaps)

    # ----------------------------------------------------
    # Slide 9: Problem Definition & Objectives
    # ----------------------------------------------------
    bullets_problem = [
        "Problem Definition:",
        "- \"To develop an efficient, lightweight, and explainable phishing website detection framework using hybrid machine learning and feature optimization techniques, capable of real-time multi-channel classification.\"",
        "Key Objectives:",
        "- Build a multi-channel engine supporting URL, email text, and SMS scans.",
        "- Implement Rough Set Theory (RSTHFS) to reduce feature space and runtime.",
        "- Train a weighted stacking ensemble of Random Forest and XGBoost.",
        "- Develop a safety-first heuristics layer to handle high-risk signals.",
        "- Integrate SHAP and Gemini API to translate math vectors into descriptive reports.",
        "- Persist scan history and profile statistics in a Neon PostgreSQL database."
    ]
    add_bullet_slide("Problem Definition & Objectives", bullets_problem)

    # ----------------------------------------------------
    # Slide 10: Proposed System Architecture
    # ----------------------------------------------------
    add_image_slide("Proposed System Architecture", "images/system_architecture_diagram.png", width_inches=8)

    # ----------------------------------------------------
    # Slide 11: Technologies Used
    # ----------------------------------------------------
    bullets_tech = [
        "Frontend Application (UI):",
        "- React.js + Vite: Light, component-driven architecture for rapid loading.",
        "- HTML5, CSS3: Custom-tailored aesthetics, visual gauge charts, and widgets.",
        "Backend Inference Engine:",
        "- Python + FastAPI: Fast, asynchronous REST API framework.",
        "- Uvicorn: High-performance ASGI web server.",
        "Database Layer:",
        "- Neon PostgreSQL (Supabase client): Storing scan results and user metrics.",
        "ML & NLP Libraries:",
        "- Scikit-learn: Random Forest, Standard Scaler.",
        "- XGBoost: High-performance gradient boosting library.",
        "- Google Generative AI (Gemini SDK) for explanation generation."
    ]
    add_bullet_slide("Technologies Used", bullets_tech)

    # ----------------------------------------------------
    # Slide 12: URL Feature Extraction
    # ----------------------------------------------------
    bullets_features = [
        "Extracts 28 structural, lexical, and trust features in 6 groups:",
        "1. Domain-based (6):",
        "- Domain name length, subdomain count, raw IP host, suspicious TLD, brand impersonation, domain age score.",
        "2. Structure (8):",
        "- URL length, @ symbol, double slash in path, special character count, digit ratio, hyphen count, dot count, query param count.",
        "3. Homograph (3):",
        "- Look-alike character presence, homograph count, Unicode/IDN presence.",
        "4. SSL Trust (3):",
        "- HTTPS scheme, SSL certificate validity, SSL age score.",
        "5. Keywords (5):",
        "- Specific keyword presence (login, secure, verify, update, urgency).",
        "6. Path/Redirect (3):",
        "- Path depth, redirect keywords, subdomain depth."
    ]
    add_bullet_slide("URL Feature Extraction", bullets_features)

    # ----------------------------------------------------
    # Slide 13: Feature Optimization (RSTHFS)
    # ----------------------------------------------------
    bullets_rsthfs = [
        "Concept & Workflow:",
        "- Reduces feature space by identifying a \"minimal reduct\" of essential attributes.",
        "- Phase 1: CDF-g Ranking evaluates attributes by cumulative distribution frequency gradients.",
        "- Phase 2: Rough Set Aggregation resolves indiscernibility and dependency relations.",
        "Optimization Benefits:",
        "- Dimensionality Reduction: Reduces feature count by 68.7% (from 48 down to 15-18).",
        "- Speed Gains: Cuts model training and backend inference time by 61%.",
        "- Latency: Lowers inference latency to <50ms (average 92ms for complete loop).",
        "- Memory Footprint: Reduced by 27%, making it suitable for standard browser extensions."
    ]
    add_bullet_slide("Feature Optimization (RSTHFS)", bullets_rsthfs)

    # ----------------------------------------------------
    # Slide 14: Ensemble Classification & Heuristics
    # ----------------------------------------------------
    bullets_ensemble = [
        "Weighted Stacking Ensemble Model:",
        "- Combines Random Forest (weight = 0.45) and XGBoost (weight = 0.55).",
        "- RF reduces variance (bagging), while XGBoost reduces bias (boosting).",
        "- Merged probability yields a threat score from 0 to 100.",
        "Classification Levels:",
        "- Safe (<40), Suspicious (40-59), Dangerous (60-79), Critical (>=80).",
        "Safety-First Heuristics Override Layer:",
        "- Intercepts predictions to prevent evasion of new domains.",
        "- Unicode Homograph Found -> Forces score to Critical (>=85).",
        "- Raw IP Address Host -> Forces score to Dangerous (>=80).",
        "- Brand Impersonation in Domain -> Forces score to Suspicious (>=75)."
    ]
    add_bullet_slide("Ensemble Model & Heuristics", bullets_ensemble)

    # ----------------------------------------------------
    # Slide 15: Explainable AI Module
    # ----------------------------------------------------
    bullets_xai = [
        "Explainable AI (XAI) Architecture:",
        "- Calculates per-feature contributions using SHAP (Shapley Additive Explanations) concepts.",
        "- Identifies the top 3-4 features driving the threat score (e.g., brand spoofing, missing SSL).",
        "Gemini AI Narrative Generator:",
        "- Backend passes suspicious feature details and metrics to Gemini model.",
        "- Generates a natural, user-friendly description of the threat technique.",
        "- Suggests safe alternative domains and immediate safety actions.",
        "Rule-Based Fallback Engine:",
        "- If Gemini is unavailable, the local engine maps features to static explanation templates.",
        "- Suggests security tips depending on the detected attack type (e.g. Look-alike letters)."
    ]
    add_bullet_slide("Explainable AI (XAI)", bullets_xai)

    # ----------------------------------------------------
    # Slide 16: UML Diagrams - Use Case & Class Diagram
    # ----------------------------------------------------
    add_two_images_slide("UML - Use Case & Class Diagram", "images/use_case.png", "images/class_dia.png", 
                         "Figure: Use Case Diagram", "Figure: Class Diagram")

    # ----------------------------------------------------
    # Slide 17: Data Flow Diagrams
    # ----------------------------------------------------
    add_two_images_slide("Data Flow Diagrams (DFD Level 0 & 1)", "images/dfd-0.png", "images/dfd-1.png",
                         "Figure: DFD Level 0 (Context Diagram)", "Figure: DFD Level 1 (Major Processes)")

    # ----------------------------------------------------
    # Slide 18: UML Sequence & Activity Diagrams
    # ----------------------------------------------------
    add_two_images_slide("UML - Sequence & Activity Diagrams", "images/sequence.png", "images/activity_diagram.png",
                         "Figure: Sequence Diagram", "Figure: Activity Diagram")

    # ----------------------------------------------------
    # Slide 19: Screenshots - URL Scan
    # ----------------------------------------------------
    add_two_images_slide("Screenshots - URL Scanner Interface", "screenshots/url.png", "screenshots/url2.png",
                         "Figure: URL Scan Tab", "Figure: Threat Score & Risk Meter")

    # ----------------------------------------------------
    # Slide 20: Screenshots - Email & SMS Scan
    # ----------------------------------------------------
    add_two_images_slide("Screenshots - Email & SMS Scanner", "screenshots/email.png", "screenshots/sms.png",
                         "Figure: Email Scan Tab (URL Extraction)", "Figure: SMS Scan Tab (Urgency Analysis)")

    # ----------------------------------------------------
    # Slide 21: Screenshots - Dashboard & Education
    # ----------------------------------------------------
    add_two_images_slide("Screenshots - Dashboard & Education", "screenshots/dashboard.png", "screenshots/education.png",
                         "Figure: Scan Metrics Dashboard", "Figure: Educational Tips Library")

    # ----------------------------------------------------
    # Slide 22: Experimental Results
    # ----------------------------------------------------
    bullets_results = [
        "Classification Model Performance Comparison:",
        "- Decision Tree: 94.32% Accuracy | 93.85% Precision | 94.10% Recall",
        "- Random Forest: 97.84% Accuracy | 97.51% Precision | 97.92% Recall",
        "- XGBoost: 98.67% Accuracy | 98.42% Precision | 98.75% Recall",
        "- Proposed Ensemble (RF + XGB): 99.21% Accuracy | 99.05% Precision | 99.18% Recall",
        "Key Latency Findings:",
        "- Complete scan pipeline (extract, predict, explain) average latency: 92ms.",
        "- Model prediction latency: 0.023 seconds (23ms) per URL.",
        "Real-World Scalability:",
        "- Tested under simulated traffic using Flask API.",
        "- Maintained accuracy above 97% and latency <50ms at 500 concurrent requests/sec."
    ]
    add_bullet_slide("Experimental Results", bullets_results)

    # ----------------------------------------------------
    # Slide 23: Conclusion & Future Scope
    # ----------------------------------------------------
    bullets_conclusion = [
        "Conclusion:",
        "- Delivered PhantomShield, a lightweight, real-time, explainable phishing detector.",
        "- Proved that combining feature selection (RSTHFS) and Stacking Ensemble yields high accuracy (99.2%) and low latency (<50ms).",
        "- Bridged the gap of \"black-box\" classifiers by calculating feature contribution (SHAP) and generating user-friendly reports.",
        "Future Scope:",
        "- Integration of Deep Learning models (Convolutional Neural Networks, Graph Neural Networks) for campaign-level tracking.",
        "- Contextual, semantic analysis of web page content using Large Language Models.",
        "- Adapt model for instant messaging (WhatsApp, Telegram) and mobile application platforms."
    ]
    add_bullet_slide("Conclusion & Future Scope", bullets_conclusion)

    # ----------------------------------------------------
    # Slide 24: Thank You Slide (Add Layout 0 slide at the end)
    # ----------------------------------------------------
    slide_thank = prs.slides.add_slide(prs.slide_layouts[0])
    slide_thank.shapes.title.text = "Thank You !!!"
    set_shape_font(slide_thank.shapes.title, font_size_pt=44, bold=True, color_rgb=RGBColor(0x0F, 0x2C, 0x59))
    
    # Use placeholder index 4 for subtitle
    if len(slide_thank.placeholders) > 0:
        for placeholder in slide_thank.placeholders:
            if placeholder.placeholder_format.idx == 4:
                placeholder.text = "Questions & Answers\nDepartment of Information Technology"
                set_shape_font(placeholder, font_size_pt=20, bold=False, color_rgb=RGBColor(0x33, 0x33, 0x33))

    # Save presentation
    output_path = 'PhantomShield_Presentation.pptx'
    prs.save(output_path)
    print(f"Successfully generated presentation at {output_path}!")

if __name__ == '__main__':
    main()
